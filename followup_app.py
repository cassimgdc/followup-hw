import gc
import json
import os
import re
import sys
import textwrap
import threading
import tkinter as tk
from datetime import datetime
from itertools import islice
from tkinter import filedialog, messagebox, ttk

import matplotlib

matplotlib.use("Agg")

import matplotlib.dates as mdates
import matplotlib.pyplot as plt
import matplotlib.ticker as mticker
import numpy as np
import pandas as pd
from pandas.api.types import is_numeric_dtype, is_object_dtype, is_string_dtype
from matplotlib.backends.backend_pdf import PdfPages
from matplotlib.offsetbox import AnnotationBbox, OffsetImage
from matplotlib.patches import FancyBboxPatch
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


def safe_div(numerator: float, denominator: float, *, multiplier: float = 1.0, fallback: float = np.nan) -> float:
    try:
        denom_value = float(denominator)
    except (TypeError, ValueError):
        denom_value = 0.0

    if not np.isfinite(denom_value) or abs(denom_value) < 1e-12:
        return fallback

    try:
        num_value = float(numerator)
    except (TypeError, ValueError):
        num_value = float(np.nan)

    if not np.isfinite(num_value):
        return fallback

    return (num_value / denom_value) * multiplier


plt.rcParams.update(
    {
        "axes.edgecolor": "#bbb",
        "axes.labelcolor": "#000000",
        "xtick.color": "#000000",
        "ytick.color": "#000000",
        "font.size": 10,
        "grid.color": "#d9d9d9",
        "grid.linestyle": ":",
        "axes.grid": True,
    }
)


class FollowUpApp:
    FILTER_COLUMNS = [
        "gNodeB Name",
        "Cell Name",
        "NR Cell ID",
        "Frequency Band",
    ]
    _NULL_TOKENS = (
        "",
        "-",
        "--",
        "---",
        "NIL",
        "Nil",
        "nil",
        "NULL",
        "Null",
        "null",
        "NONE",
        "None",
        "none",
        "NaN",
        "nan",
        "N/A",
        "n/a",
        "NA",
        "na",
    )
    _EVAL_GLOBALS: dict[str, object] = {
        "np": np,
        "pd": pd,
        "min": min,
        "max": max,
        "abs": abs,
        "round": round,
        "SAFE_DIV": safe_div,
    }
    def __init__(self, master):
        self.master = master
        self.master.title("5G Follow-Up Generator - Huawei/VIVO")
        self.master.protocol("WM_DELETE_WINDOW", self.on_close)
        self.csv_paths = []
        self.tasklines = []
        self.df: pd.DataFrame | None = None
        self._is_generating = False
        self._app_dir = self._resolve_app_dir()
        self._resource_dirs = self._build_resource_directories()
        self._header_row_cache: dict[str, int] = {}
        self.build_interface()

    @staticmethod
    def _format_filter_entry(value) -> str:
        if pd.isna(value):
            return ""
        if isinstance(value, (int, np.integer)):
            return str(value)
        if isinstance(value, float):
            if np.isfinite(value) and float(value).is_integer():
                return str(int(value))
            return str(value)
        return str(value)

    def _sort_filter_values(self, values: set[str]) -> list[str]:
        def sort_key(val: str) -> tuple[int, object]:
            if val is None or val == "":
                return (2, "")
            try:
                numeric_val = float(val)
                if float(numeric_val).is_integer():
                    numeric_val = int(numeric_val)
                return (0, numeric_val)
            except (TypeError, ValueError):
                return (1, str(val).lower())

        return sorted({self._format_filter_entry(v) for v in values if v is not None}, key=sort_key)

    def _detect_data_start_line(self, path: str) -> int:
        cached = self._header_row_cache.get(path)
        if cached is not None:
            return cached

        header_idx = 0
        with open(path, "r", encoding="utf-8", errors="ignore") as handle:
            for idx, line in enumerate(handle):
                prefix = line.strip().lower()
                if prefix.startswith("time") or prefix.startswith("day"):
                    header_idx = idx
                    break
        self._header_row_cache[path] = header_idx
        return header_idx

    def _collect_filter_selections(self) -> dict[str, list[str]]:
        selections: dict[str, list[str]] = {}
        for col, lb in self.filtros.items():
            selected_values = [lb.get(i) for i in lb.curselection()]
            if selected_values:
                selections[col] = selected_values
        return selections

    def _build_filter_slug(self, selections: dict[str, list[str]]) -> str:
        tokens: list[str] = []
        seen: set[str] = set()
        ordered_keys = ["gNodeB Name", "Cell Name", "NR Cell ID", "Frequency Band"]
        for key in ordered_keys:
            for value in selections.get(key, []):
                slug = self._slugify(value)
                if slug and slug not in seen:
                    tokens.append(slug)
                    seen.add(slug)
        return "_".join(tokens)

    def _compose_base_folder(self, followup_slug: str, date_tag: str) -> str:
        parts = ["FollowUp_5G"]
        if followup_slug:
            parts.append(followup_slug)
        parts.append(date_tag)
        return "_".join(parts)

    def _compose_file_basename(self, followup_slug: str, filter_slug: str, date_tag: str) -> str:
        parts = ["FollowUp_5G"]
        if followup_slug:
            parts.append(followup_slug)
        if filter_slug:
            parts.append(filter_slug)
        parts.append(date_tag)
        return "_".join(parts)

    def _resolve_app_dir(self) -> str:
        meipass = getattr(sys, "_MEIPASS", None)
        if meipass:
            return meipass
        return os.path.dirname(os.path.abspath(__file__))

    def _build_resource_directories(self) -> list[str]:
        dirs: list[str] = []

        exe_path = getattr(sys, "executable", None)
        if exe_path:
            exe_dir = os.path.dirname(os.path.abspath(exe_path))
            if exe_dir and exe_dir not in dirs:
                dirs.append(exe_dir)

        cwd = os.getcwd()
        if cwd and cwd not in dirs:
            dirs.append(cwd)

        if self._app_dir not in dirs:
            dirs.append(self._app_dir)

        return dirs

    def _find_resource(self, filename: str) -> str | None:
        for base_dir in self._resource_dirs:
            direct_path = os.path.join(base_dir, filename)
            if os.path.exists(direct_path):
                return direct_path

            assets_path = os.path.join(base_dir, "assets", filename)
            if os.path.exists(assets_path):
                return assets_path

        return None

    def validate_date_entry(self, event):
        widget = event.widget
        try:
            datetime.strptime(widget.get(), "%Y-%m-%d")
        except ValueError:
            messagebox.showerror("Data inválida", "Use o formato AAAA-MM-DD.")
            widget.focus_set()

    def validate_datetime_entry(self, event):
        widget = event.widget
        try:
            datetime.strptime(widget.get(), "%Y-%m-%d %H:%M")
        except ValueError:
            messagebox.showerror("Data e hora inválidas", "Use o formato AAAA-MM-DD HH:MM.")
            widget.focus_set()

    def build_interface(self):
        frm = ttk.Frame(self.master, padding=10)
        frm.pack(fill="both", expand=True)

        ttk.Label(frm, text="Arquivos CSV (5G1 e 5G2):").pack(anchor="w")
        ttk.Button(frm, text="Selecionar Arquivos", command=self.select_csvs).pack(pady=4)
        self.files_label = ttk.Label(frm, text="Nenhum arquivo selecionado.")
        self.files_label.pack(anchor="w")

        self.filtros: dict[str, tk.Listbox] = {}
        for col in ["gNodeB Name", "Cell Name", "NR Cell ID", "Frequency Band"]:
            ttk.Label(frm, text=f"Filtro: {col}").pack(anchor="w")
            lb = tk.Listbox(frm, selectmode="multiple", height=4, exportselection=False)
            lb.pack(fill="x", pady=2)
            self.filtros[col] = lb

        ttk.Label(frm, text="Intervalo de Datas (AAAA-MM-DD):").pack(anchor="w", pady=(10, 0))
        dtf = ttk.Frame(frm)
        dtf.pack(fill="x")
        self.date_ini = ttk.Entry(dtf, width=12)
        self.date_ini.pack(side="left", padx=5)
        self.date_ini.bind("<FocusOut>", self.validate_date_entry)
        self.date_fim = ttk.Entry(dtf, width=12)
        self.date_fim.pack(side="left", padx=5)
        self.date_fim.bind("<FocusOut>", self.validate_date_entry)

        ttk.Label(frm, text="Filtro: Hora do dia").pack(anchor="w")
        self.hour_filter = tk.Listbox(frm, selectmode="multiple", height=6, exportselection=False)
        for h in range(24):
            self.hour_filter.insert("end", f"{h:02}:00")
        self.hour_filter.pack(fill="x", pady=(0, 6))

        ttk.Label(frm, text="Tasklines (evento + horário):").pack(anchor="w")
        taskf = ttk.Frame(frm)
        taskf.pack(fill="x")
        self.task_entry_time = ttk.Entry(taskf)
        self.task_entry_time.pack(side="left", padx=2)
        self.task_entry_time.bind("<FocusOut>", self.validate_datetime_entry)
        self.task_entry_label = ttk.Entry(taskf)
        self.task_entry_label.insert(0, "Ativação")
        self.task_entry_label.pack(side="left", padx=2)
        ttk.Button(taskf, text="+", width=3, command=self.add_taskline).pack(side="left", padx=2)

        self.task_listbox = tk.Listbox(frm, height=3)
        self.task_listbox.pack(fill="x", pady=4)
        self.task_listbox.bind("<Double-Button-1>", self.remove_taskline)

        ttk.Button(frm, text="Limpar Filtros", command=self.reset_filters).pack(pady=(0, 10))

        self.kpi_title_entry = ttk.Entry(frm)
        self.kpi_title_entry.insert(0, "Nome do Follow-Up (opcional)")
        self.kpi_title_entry.pack(fill="x", pady=(0, 10))

        exportf = ttk.Frame(frm)
        exportf.pack(fill="x", pady=10)
        self.export_pdf = tk.BooleanVar(value=True)
        self.export_png = tk.BooleanVar(value=False)
        self.export_ppt = tk.BooleanVar(value=False)
        ttk.Checkbutton(exportf, text="Exportar PDF", variable=self.export_pdf).pack(
            side="left", padx=5
        )
        ttk.Checkbutton(exportf, text="Exportar PNG", variable=self.export_png).pack(
            side="left", padx=5
        )
        ttk.Checkbutton(exportf, text="Exportar PPTX", variable=self.export_ppt).pack(
            side="left", padx=5
        )

        ttk.Label(exportf, text="Gráficos por página:").pack(side="left", padx=(20, 5))
        self.page_layout = ttk.Combobox(exportf, values=[4, 6], width=5)
        self.page_layout.set(4)
        self.page_layout.pack(side="left")

        self.high_definition = tk.BooleanVar(value=False)
        ttk.Checkbutton(
            frm,
            text="Gerar em alta definição",
            variable=self.high_definition,
        ).pack(anchor="w", pady=(2, 0))

        ttk.Button(frm, text="Gerar Relatório", command=self.generate_report).pack(pady=10)

        self.progress_var = tk.DoubleVar(value=0.0)
        self.progress_frame = ttk.Frame(frm)
        self.progress_label = ttk.Label(self.progress_frame, text="")
        self.progress_label.pack(anchor="w")
        self.progress_bar = ttk.Progressbar(
            self.progress_frame, maximum=100, variable=self.progress_var
        )
        self.progress_bar.pack(fill="x", pady=(2, 0))
        self.progress_frame.pack(fill="x")
        self.progress_frame.pack_forget()

    def reset_filters(self):
        for lb in self.filtros.values():
            lb.selection_clear(0, "end")
        self.hour_filter.selection_clear(0, "end")

    def select_csvs(self):
        paths = filedialog.askopenfilenames(filetypes=[("CSV Files", "*.csv")])
        if not paths:
            return

        self.csv_paths = list(paths)
        df = self.merge_csvs()
        self.df = df
        fnames = [os.path.basename(p) for p in self.csv_paths]
        self.files_label.config(text="Selecionados: " + ", ".join(fnames))

        for col, lb in self.filtros.items():
            lb.delete(0, "end")
            if col in df.columns:
                options = self._sort_filter_values(
                    {self._format_filter_entry(value) for value in df[col].dropna().unique()}
                )
            else:
                options = []
            for display_value in options:
                lb.insert("end", display_value)

        min_dt = df["DATETIME"].dropna().min() if "DATETIME" in df.columns else None
        max_dt = df["DATETIME"].dropna().max() if "DATETIME" in df.columns else None

        if pd.notna(min_dt) and pd.notna(max_dt):
            self.date_ini.delete(0, "end")
            self.date_ini.insert(0, min_dt.strftime("%Y-%m-%d"))
            self.date_fim.delete(0, "end")
            self.date_fim.insert(0, max_dt.strftime("%Y-%m-%d"))
            self.task_entry_time.delete(0, "end")
            self.task_entry_time.insert(0, min_dt.strftime("%Y-%m-%d %H:00"))

    def add_taskline(self):
        time = self.task_entry_time.get()
        label = self.task_entry_label.get()
        if time:
            try:
                ts = pd.to_datetime(time)
            except ValueError:
                messagebox.showerror("Data inválida", "Use o formato AAAA-MM-DD HH:MM.")
                return
            self.tasklines.append((ts, label))
            self.task_listbox.insert("end", f"{time} - {label}")

    def remove_taskline(self, event):
        sel = self.task_listbox.curselection()
        if sel:
            self.tasklines.pop(sel[0])
            self.task_listbox.delete(sel[0])

    def on_close(self):
        plt.close("all")
        self.master.quit()
        self.master.destroy()

    def aplicar_filtros(
        self,
        df: pd.DataFrame,
        *,
        selections: dict[str, list[str]] | None = None,
        date_range: tuple[str | None, str | None] | None = None,
        selected_hours: list[str] | None = None,
    ) -> pd.DataFrame:
        if df is None:
            return pd.DataFrame()

        filtered = df
        if "DATETIME" in filtered.columns:
            filtered = filtered[filtered["DATETIME"].notna()]

        if selections is None:
            filter_selections: dict[str, list[str]] = {
                col: [lb.get(i) for i in lb.curselection()]
                for col, lb in self.filtros.items()
            }
        else:
            filter_selections = {col: list(values) for col, values in selections.items()}

        for col, values in filter_selections.items():
            if not values or col not in filtered.columns:
                continue
            normalized_selections = {item for item in values}
            column_values = filtered[col].apply(self._format_filter_entry)
            mask = column_values.isin(normalized_selections)
            filtered = filtered[mask]

        if date_range is None:
            start_value = self.date_ini.get()
            end_value = self.date_fim.get()
        else:
            start_value, end_value = date_range

        ini = pd.to_datetime(start_value, errors="coerce")
        fim = pd.to_datetime(end_value, errors="coerce")
        if pd.notna(ini) and pd.notna(fim) and "DATETIME" in filtered.columns:
            filtered = filtered[(filtered["DATETIME"] >= ini) & (filtered["DATETIME"] <= fim)]

        if selected_hours is None:
            hour_strings = [self.hour_filter.get(i) for i in self.hour_filter.curselection()]
        else:
            hour_strings = list(selected_hours)

        if hour_strings and "DATETIME" in filtered.columns:
            try:
                hour_values = {int(item.split(":", 1)[0]) for item in hour_strings}
            except ValueError:
                hour_values = set()
            if hour_values:
                filtered = filtered[filtered["DATETIME"].dt.hour.isin(hour_values)]

        return filtered

    def clean_csv(self, path: str) -> pd.DataFrame:
        skiprows = self._detect_data_start_line(path)
        try:
            return pd.read_csv(
                path,
                skiprows=skiprows,
                skipfooter=1,
                engine="python",
            )
        except ValueError:
            return pd.read_csv(
                path,
                skiprows=skiprows,
                engine="python",
            )

    def merge_csvs(self) -> pd.DataFrame:
        dfs = [self.clean_csv(path) for path in self.csv_paths]
        dfs = [df for df in dfs if not df.empty]
        if not dfs:
            return pd.DataFrame()

        df = pd.concat(dfs, axis=1)
        df = df.loc[:, ~df.columns.duplicated()].copy()

        if "Time" in df.columns:
            df["DATETIME"] = pd.to_datetime(df["Time"], errors="coerce")
        elif "Day" in df.columns:
            df["DATETIME"] = pd.to_datetime(df["Day"], errors="coerce")

        if "DATETIME" in df.columns:
            df.sort_values(by="DATETIME", inplace=True, ignore_index=True)

        return df

    def load_kpi_formulas(self) -> dict:
        path = self._find_resource("kpi_formulas.json")
        if not path:
            return {}

        try:
            with open(path, "r", encoding="utf-8") as f:
                return json.load(f)
        except FileNotFoundError:
            return {}
        except json.JSONDecodeError as exc:
            messagebox.showerror(
                "Erro",
                f"Não foi possível ler kpi_formulas.json: {exc}",
            )
            return {}

    _SOURCE_PATTERN = re.compile(r"'[^']+'\[([^\]]+)\]")
    _DATAFRAME_PATTERN = re.compile(r"df\[(?:'([^']+)'|\"([^\"]+)\")\]")

    def _replace_source_refs(self, formula: str) -> str:
        def repl(match: re.Match) -> str:
            column = match.group(1)
            return f"df[{repr(column)}]"

        return self._SOURCE_PATTERN.sub(repl, formula)

    def _extract_formula_columns(self, formula: str) -> set[str]:
        columns: set[str] = set()
        for match in self._SOURCE_PATTERN.finditer(formula):
            columns.add(match.group(1))

        replaced = self._replace_source_refs(formula)
        for match in self._DATAFRAME_PATTERN.finditer(replaced):
            column = match.group(1) or match.group(2)
            if column:
                columns.add(column)

        columns.discard("DATETIME")
        return columns

    def _sanitize_numeric_series(self, series: pd.Series) -> pd.Series:
        if series.empty:
            return series.astype(np.float32, copy=False)

        if is_numeric_dtype(series):
            numeric = pd.to_numeric(series, errors="coerce")
            return numeric.astype(np.float32, copy=False)

        working = series.replace(self._NULL_TOKENS, np.nan)

        if is_object_dtype(working) or is_string_dtype(working):
            working = working.astype(str).str.strip()
            working = working.replace(self._NULL_TOKENS, np.nan)
            working = working.replace(r"^$", np.nan, regex=True)
            working = working.str.replace(",", ".", regex=False)

        numeric = pd.to_numeric(working, errors="coerce")
        return numeric.astype(np.float32, copy=False)

    def _sanitize_numeric_columns(self, df: pd.DataFrame, columns: set[str]) -> None:
        if not columns:
            return

        for column in columns:
            if column not in df.columns:
                continue

            series = df[column]
            if is_numeric_dtype(series) and series.dtype == np.float32:
                continue

            df[column] = self._sanitize_numeric_series(series)

    def _normalize_formula_output(self, value, kpi_name: str) -> pd.Series | None:
        if isinstance(value, pd.Series):
            numeric = pd.to_numeric(value, errors="coerce")
            numeric = numeric.dropna()
            return numeric if not numeric.empty else None

        if isinstance(value, pd.DataFrame):
            numeric_df = value.select_dtypes(include=[np.number])
            if numeric_df.empty:
                return None
            collapsed = numeric_df.sum(axis=0)
            collapsed = pd.to_numeric(collapsed, errors="coerce").dropna()
            return collapsed if not collapsed.empty else None

        if isinstance(value, dict):
            series = pd.Series(value, dtype="float64")
            series = pd.to_numeric(series, errors="coerce").dropna()
            return series if not series.empty else None

        if isinstance(value, (list, tuple, np.ndarray)):
            if len(value) == 0:
                return None
            series = pd.Series(value, dtype="float64")
            series.index = [f"{kpi_name} {idx}" for idx in range(1, len(series) + 1)]
            series = pd.to_numeric(series, errors="coerce").dropna()
            return series if not series.empty else None

        try:
            numeric_value = float(value)
        except (TypeError, ValueError):
            return None

        if not np.isfinite(numeric_value):
            return None

        return pd.Series({kpi_name: numeric_value}, dtype="float64")

    def _evaluate_formula(
        self,
        grouped: "pd.core.groupby.generic.DataFrameGroupBy",
        compiled_formula,
        kpi_name: str,
    ) -> pd.DataFrame:
        if grouped.ngroups == 0:
            return pd.DataFrame()

        rows: list[pd.Series] = []
        index_values: list[pd.Timestamp] = []
        for timestamp, group in grouped:
            local_context = {"df": group}
            try:
                value = eval(compiled_formula, self._EVAL_GLOBALS, local_context)
            except Exception as exc:
                raise ValueError(
                    f"Erro na fórmula do KPI '{kpi_name}': {exc}"
                ) from exc

            series = self._normalize_formula_output(value, kpi_name)
            if series is None:
                continue

            rows.append(series)
            index_values.append(pd.to_datetime(timestamp))

        if not rows:
            return pd.DataFrame()

        df_result = pd.DataFrame(rows, index=index_values)
        df_result.sort_index(inplace=True)
        df_result.replace([np.inf, -np.inf], np.nan, inplace=True)
        df_result.dropna(how="all", inplace=True)
        if df_result.empty:
            return df_result

        return df_result.astype(np.float32, copy=False)

    @staticmethod
    def _chunked(iterable, size):
        iterator = iter(iterable)
        while True:
            chunk = list(islice(iterator, size))
            if not chunk:
                break
            yield chunk

    def _is_percentage_metric(self, formula: str) -> bool:
        formula = formula or ""
        normalized = re.sub(r"\s+", "", formula)
        if re.search(r"(?i)multiplier=100(?!\d)", normalized):
            return True

        # Captures expressions such as *100, *100.0 or 100*(...) without matching 1000, 10000...
        if re.search(r"\*100(?!\d)", normalized):
            return True

        if re.search(r"(?<!\d)100(?:\.0+)?\*", normalized):
            return True

        return False

    def _load_brand_assets(self) -> dict[str, str]:
        expected = {
            "huawei": self._find_resource("huawei.png"),
            "vivo": self._find_resource("vivo.png"),
            "background": self._find_resource("mapa-fundo.png"),
        }

        missing = [name for name, path in expected.items() if not path]
        if missing:
            missing_names = ", ".join(missing)
            search_locations = "\n".join(self._resource_dirs)
            raise FileNotFoundError(
                "Arquivos de branding ausentes: "
                f"{missing_names}. Posicione os arquivos na mesma pasta do executável"
                " ou em uma subpasta 'assets/'.\nPastas verificadas:\n"
                f"{search_locations}"
            )

        return {key: path for key, path in expected.items() if path}

    def _gather_report_metadata(
        self,
        df_filt: pd.DataFrame,
        *,
        filter_selections: dict[str, list[str]] | None = None,
        date_ini: str | None = None,
        date_fim: str | None = None,
        selected_hours: list[str] | None = None,
        tasklines: list[tuple[pd.Timestamp, str]] | None = None,
        report_name_entry: str | None = None,
    ) -> dict:
        placeholder_name = "Nome do Follow-Up (opcional)"

        if filter_selections is None:
            filter_selections = self._collect_filter_selections()
        filters_summary: dict[str, str] = {}
        for col in self.filtros.keys():
            selections = filter_selections.get(col, [])
            filters_summary[col] = ", ".join(selections) if selections else "All"
        filter_slug = self._build_filter_slug(filter_selections)

        start_value = (date_ini or self.date_ini.get()).strip()
        end_value = (date_fim or self.date_fim.get()).strip()
        if start_value and end_value:
            date_range = f"{start_value} to {end_value}"
        elif start_value or end_value:
            date_range = start_value or end_value
        else:
            date_range = "All"

        if selected_hours is None:
            hours_snapshot = [self.hour_filter.get(i) for i in self.hour_filter.curselection()]
        else:
            hours_snapshot = list(selected_hours)
        hour_summary = ", ".join(hours_snapshot) if hours_snapshot else "All"

        if tasklines is None:
            tasklines_iter = list(self.tasklines)
        else:
            tasklines_iter = list(tasklines)
        formatted_tasklines: list[str] = []
        for idx, (ts, label) in enumerate(tasklines_iter, start=1):
            ts_fmt = ts.strftime("%Y-%m-%d %H:%M")
            formatted_tasklines.append(f"{idx}. {label} ({ts_fmt})")

        if report_name_entry is None:
            report_name_entry = self.kpi_title_entry.get().strip()
        report_name = (
            report_name_entry
            if report_name_entry and report_name_entry != placeholder_name
            else ""
        )

        return {
            "report_name": report_name,
            "filters": filters_summary,
            "date_range": date_range,
            "hours": hour_summary,
            "tasklines": formatted_tasklines,
            "total_registros": len(df_filt),
            "filter_slug": filter_slug,
        }

    def _slugify(self, value: str) -> str:
        if not value:
            return ""
        normalized = re.sub(r"[^\w-]+", "_", value, flags=re.UNICODE)
        normalized = re.sub(r"_+", "_", normalized)
        return normalized.strip("_")

    def _normalize_report_name(self, raw_value: str) -> str:
        placeholder_name = "Nome do Follow-Up (opcional)"
        cleaned = (raw_value or "").strip()
        if not cleaned or cleaned == placeholder_name:
            return ""
        return self._slugify(cleaned)

    def _add_cover_annotation(self, ax, image_path: str, xy: tuple[float, float], zoom: float):
        if not os.path.exists(image_path):
            return
        image = plt.imread(image_path)
        image_box = OffsetImage(image, zoom=zoom)
        box = AnnotationBbox(image_box, xy, frameon=False)
        ax.add_artist(box)

    def _build_summary_lines(self, metadata: dict) -> list[str]:
        lines: list[str] = []
        if metadata["report_name"]:
            lines.append(metadata["report_name"])
        lines.append(f"Date: {metadata['date_range']}")
        lines.append(f"Hour: {metadata['hours']}")
        label_map = {
            "gNodeB Name": "gNodeB",
            "Cell Name": "Cell",
            "NR Cell ID": "NR Cell ID",
            "Frequency Band": "Frequency Band",
        }
        for key, label in label_map.items():
            value = metadata["filters"].get(key)
            if value is not None:
                lines.append(f"{label}: {value}")
        if metadata["tasklines"]:
            lines.append("Task Lines:")
            lines.extend(metadata["tasklines"])
        return lines

    def _shorten_lines(self, lines: list[str], width: int = 85) -> list[str]:
        shortened: list[str] = []
        for line in lines:
            if not line:
                shortened.append("")
                continue
            shortened.append(textwrap.shorten(line, width=width, placeholder="..."))
        return shortened

    def _build_pdf_cover(
        self,
        pdf: PdfPages,
        metadata: dict,
        logos: dict[str, str],
        page_size: tuple[float, float],
        dpi: int,
    ):
        fig, ax = plt.subplots(figsize=page_size)
        ax.axis("off")
        fig.patch.set_facecolor("#ffffff")
        ax.set_xlim(0, 1)
        ax.set_ylim(0, 1)

        background = plt.imread(logos["background"])
        image_artist = ax.imshow(
            background,
            extent=[0.42, 1.08, -0.03, 1.03],
            aspect="auto",
            alpha=1.0,
            zorder=0,
            origin="upper",
        )
        image_artist.set_clip_on(False)

        self._add_cover_annotation(ax, logos["huawei"], (0.08, 0.94), 0.48)
        self._add_cover_annotation(ax, logos["vivo"], (0.33, 0.94), 0.48)

        ax.text(
            0.12,
            0.78,
            "FOLLOW UP",
            fontsize=40,
            fontweight="bold",
            color="#c3002f",
            va="top",
            ha="left",
            fontfamily="DejaVu Sans",
        )
        ax.text(
            0.12,
            0.64,
            "5G",
            fontsize=60,
            fontweight="bold",
            color="#c3002f",
            va="top",
            ha="left",
            fontfamily="DejaVu Sans",
        )

        summary_lines = self._shorten_lines(self._build_summary_lines(metadata))
        max_chars = max((len(line) for line in summary_lines), default=0)
        line_count = max(len(summary_lines), 1)

        box_width = min(0.52, max(0.27, 0.0095 * max_chars + 0.13))
        box_height = min(0.48, max(0.035 + 0.030 * line_count, 0.14))
        box_x = 0.12
        box_y = 0.04

        box = FancyBboxPatch(
            (box_x, box_y),
            box_width,
            box_height,
            boxstyle="round,pad=0.08",
            linewidth=1.6,
            edgecolor="#c3002f",
            facecolor="#fff7f7",
        )
        ax.add_patch(box)

        ax.text(
            box_x + 0.02,
            box_y + box_height - 0.02,
            "\n".join(summary_lines),
            fontsize=9,
            color="#4b1a20",
            va="top",
            ha="left",
            linespacing=1.15,
            fontfamily="DejaVu Sans Mono",
        )

        pdf.savefig(fig, dpi=dpi)
        plt.close(fig)

    def _add_pdf_header(self, fig: plt.Figure, metadata: dict, logos: dict[str, str]):
        header_height = 0.1
        fig.subplots_adjust(top=0.82)
        logo_alpha = 0.35

        left_ax = fig.add_axes([0.02, 0.86, 0.12, header_height])
        left_ax.axis("off")
        left_ax.set_zorder(-1)
        if os.path.exists(logos["huawei"]):
            left_ax.patch.set_alpha(0.0)
            left_ax.imshow(plt.imread(logos["huawei"]), alpha=logo_alpha, zorder=-1)

        right_ax = fig.add_axes([0.86, 0.86, 0.12, header_height])
        right_ax.axis("off")
        right_ax.set_zorder(-1)
        if os.path.exists(logos["vivo"]):
            right_ax.patch.set_alpha(0.0)
            right_ax.imshow(plt.imread(logos["vivo"]), alpha=logo_alpha, zorder=-1)

        label_map = {
            "gNodeB Name": "gNodeB",
            "Cell Name": "Cell",
            "NR Cell ID": "NR Cell ID",
            "Frequency Band": "Frequency Band",
        }
        filters_summary = [f"Hour: {metadata['hours']}"]
        for key in label_map:
            value = metadata["filters"].get(key)
            if value is not None:
                filters_summary.append(f"{label_map[key]}: {value}")

        header_text = "\n".join(filters_summary)
        fig.text(
            0.5,
            0.94,
            header_text,
            ha="center",
            va="top",
            fontsize=6,
            color="#333333",
        )

    def _add_ppt_cover(self, prs: Presentation, metadata: dict, logos: dict[str, str]):
        blank = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank)

        if os.path.exists(logos["background"]):
            slide.shapes.add_picture(
                logos["background"],
                Inches(3.2),
                Inches(0.2),
                width=prs.slide_width - Inches(3.4),
            )

        if os.path.exists(logos["huawei"]):
            slide.shapes.add_picture(logos["huawei"], Inches(1.4), Inches(0.2), height=Inches(1.65))
        if os.path.exists(logos["vivo"]):
            slide.shapes.add_picture(logos["vivo"], Inches(3.6), Inches(0.26), height=Inches(1.45))

        title_box = slide.shapes.add_textbox(Inches(1.2), Inches(1.45), Inches(3.8), Inches(2.4))
        title_frame = title_box.text_frame
        title_frame.clear()
        title_frame.text = "FOLLOW UP"
        title_frame.paragraphs[0].font.size = Pt(38)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(195, 0, 47)
        title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        title_frame.paragraphs[0].font.name = "Arial"
        para = title_frame.add_paragraph()
        para.text = "5G"
        para.font.size = Pt(64)
        para.font.bold = True
        para.font.color.rgb = RGBColor(195, 0, 47)
        para.alignment = PP_ALIGN.LEFT
        para.font.name = "Arial"

        lines = self._shorten_lines(self._build_summary_lines(metadata))
        max_chars = max((len(line) for line in lines), default=0)
        line_count = max(len(lines), 1)
        box_width_in = min(6.0, max(3.4, 0.09 * max_chars + 1.8))
        box_height_in = min(5.0, max(2.0, 0.32 * line_count))

        summary_shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(1.2),
            prs.slide_height - Inches(2.35),
            Inches(box_width_in),
            Inches(box_height_in),
        )
        fill = summary_shape.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 245, 245)
        line = summary_shape.line
        line.color.rgb = RGBColor(195, 0, 47)
        line.width = Pt(2)

        summary_frame = summary_shape.text_frame
        summary_frame.clear()
        if not lines:
            summary_frame.text = ""
        else:
            for idx, text in enumerate(lines):
                if idx == 0:
                    paragraph = summary_frame.paragraphs[0]
                    paragraph.text = text
                else:
                    paragraph = summary_frame.add_paragraph()
                    paragraph.text = text

                if idx == 0 and metadata["report_name"]:
                    paragraph.font.size = Pt(18)
                    paragraph.font.bold = True
                else:
                    paragraph.font.size = Pt(9.5)
                    paragraph.font.bold = False
                paragraph.font.color.rgb = RGBColor(75, 26, 32)
                paragraph.font.name = "Arial"

    def _render_chart(
        self,
        ax: plt.Axes,
        data: pd.DataFrame,
        kpi_name: str,
        is_percent: bool,
        chart_type: str,
        *,
        style_scale: float = 1.0,
        tasklines: list[tuple[pd.Timestamp, str]] | None = None,
    ) -> None:
        fig = ax.figure
        fig.patch.set_facecolor("#ffffff")

        chart_kind = (chart_type or "line").lower()
        base_facecolor = "#faf3f4" if chart_kind == "area" else "#ffffff"
        ax.set_facecolor(base_facecolor)

        for spine in ax.spines.values():
            spine.set_visible(False)
        line_color = "#a23746"
        fill_color = "#c45562"
        linewidth = 1.1 * max(0.6, style_scale)

        plot_df = data.sort_index().astype(float)
        if plot_df.empty:
            return

        if is_percent:
            finite_values = plot_df.replace([np.inf, -np.inf], np.nan).stack(dropna=True)
            if not finite_values.empty and finite_values.max() <= 1.5:
                plot_df = plot_df * 100.0

        x_values = plot_df.index
        legend_labels: list[str] = []

        if chart_kind == "stacked_bar":
            date_numbers = mdates.date2num(x_values)
            if len(date_numbers) > 1:
                spacing = np.diff(np.sort(date_numbers))
                bar_width = float(np.min(spacing)) * 0.8
            else:
                bar_width = 0.25
            bottoms = np.zeros(len(plot_df))
            palette = [
                "#a23746",
                "#c45562",
                "#de6a73",
                "#ee8e92",
                "#f5b5b7",
                "#fbe0df",
            ]
            for idx, column in enumerate(plot_df.columns):
                values = plot_df[column].fillna(0).to_numpy()
                color = palette[idx % len(palette)]
                ax.bar(
                    x_values,
                    values,
                    width=bar_width,
                    bottom=bottoms,
                    color=color,
                    edgecolor="#ffffff",
                    linewidth=0.4,
                )
                bottoms = bottoms + values
                legend_labels.append(str(column))
            ax.grid(False)
        elif chart_kind == "bar":
            date_numbers = mdates.date2num(x_values)
            if len(date_numbers) > 1:
                spacing = np.diff(np.sort(date_numbers))
                bar_width = float(np.min(spacing)) * 0.7
            else:
                bar_width = 0.18
            column = plot_df.columns[0]
            ax.bar(
                x_values,
                plot_df[column].to_numpy(),
                width=bar_width,
                color="#cc5463",
                alpha=0.9,
                edgecolor="#8c1f33",
            )
            legend_labels.append(str(column))
            ax.grid(False)
        else:
            column = plot_df.columns[0]
            ax.plot(
                x_values,
                plot_df[column].to_numpy(),
                linewidth=linewidth,
                color=line_color,
                zorder=5,
            )
            if chart_kind == "area":
                ax.fill_between(
                    x_values,
                    plot_df[column].to_numpy(),
                    color=fill_color,
                    alpha=0.2,
                    zorder=4,
                )
            legend_labels.append(str(column))

        if chart_kind not in {"stacked_bar", "bar"}:
            ax.grid(True)

        ymin, ymax = ax.get_ylim()
        if np.isfinite(ymin) and np.isfinite(ymax):
            padding = (ymax - ymin) * 0.1 if ymax != ymin else max(abs(ymax), 1.0) * 0.1
        else:
            padding = 1.0
        ax.set_ylim(ymin - padding * 0.12, ymax + padding)

        for ts, label in tasklines or []:
            if x_values.min() <= ts <= x_values.max():
                ax.axvline(ts, color="#c3002f", linestyle="--", linewidth=1.0, zorder=2)
                ax.text(
                    ts,
                    ax.get_ylim()[1] - padding * 0.25,
                    label,
                    color="#7a0018",
                    rotation=90,
                    ha="right",
                    va="top",
                    fontsize=max(6, 7.2 * style_scale),
                    bbox={
                        "facecolor": "#ffffff",
                        "edgecolor": "none",
                        "alpha": 0.7,
                        "pad": 1,
                    },
                    zorder=1,
                )

        ax.set_xlabel("", color="#000000")
        ax.set_ylabel("", color="#000000")
        ax.tick_params(axis="both", colors="#000000", labelsize=max(7, 9 * style_scale))

        if is_percent:
            ax.yaxis.set_major_formatter(mticker.PercentFormatter(xmax=100, decimals=1))

        ax.xaxis.set_major_locator(mdates.AutoDateLocator())
        ax.xaxis.set_major_formatter(mdates.DateFormatter("%d/%m %H:%M"))
        ax.margins(x=0.01)

        wrapped_title = textwrap.fill(kpi_name, width=40 if style_scale >= 1.0 else 34)
        title_font = max(8, 11 * style_scale)
        ax.set_title(
            wrapped_title,
            fontweight="bold",
            fontsize=title_font,
            color="#6c0b14",
            pad=20 * style_scale,
        )
        title = ax.title
        if chart_kind in {"stacked_bar", "bar"}:
            title.set_y(1.14 if style_scale >= 1.0 else 1.12)
        else:
            title.set_y(1.04 if style_scale < 1.0 else 1.05)

        if legend_labels and chart_kind in {"stacked_bar", "bar"} and len(plot_df.columns) > 1:
            legend_anchor = 1.06 if style_scale >= 0.9 else 1.04
            ax.legend(
                legend_labels,
                loc="upper center",
                bbox_to_anchor=(0.5, legend_anchor),
                ncol=min(len(legend_labels), 4),
                frameon=False,
                fontsize=max(6, 6.2 * style_scale),
            )

        for label in ax.get_xticklabels():
            label.set_rotation(24)
            label.set_horizontalalignment("right")

    def _render_placeholder_chart(
        self,
        ax: plt.Axes,
        kpi_name: str,
        *,
        style_scale: float = 1.0,
    ) -> None:
        fig = ax.figure
        fig.patch.set_facecolor("#ffffff")
        ax.set_facecolor("#ffffff")
        for spine in ax.spines.values():
            spine.set_visible(False)
        ax.grid(False)
        ax.set_xticks([])
        ax.set_yticks([])
        ax.set_xlim(0, 1)
        ax.set_ylim(0, 1)

        wrapped_title = textwrap.fill(kpi_name, width=40 if style_scale >= 1.0 else 34)
        title_font = max(8, 11 * style_scale)
        ax.set_title(
            wrapped_title,
            fontweight="bold",
            fontsize=title_font,
            color="#6c0b14",
            pad=20 * style_scale,
        )
        ax.text(
            0.5,
            0.5,
            "Sem dados disponíveis",
            ha="center",
            va="center",
            fontsize=max(9, 11 * style_scale),
            color="#7a0018",
            fontweight="bold",
        )

    def generate_report(self):
        if self._is_generating:
            messagebox.showinfo(
                "Processando",
                "Um relatório já está sendo gerado. Aguarde a finalização antes de iniciar outro.",
            )
            return

        if self.df is None:
            messagebox.showerror("Erro", "Importe os CSVs antes.")
            return

        should_export_png = self.export_png.get()
        should_export_pdf = self.export_pdf.get()
        should_export_ppt = self.export_ppt.get()
        if not any([should_export_png, should_export_pdf, should_export_ppt]):
            messagebox.showerror(
                "Formato não selecionado",
                "Selecione pelo menos um formato de exportação (PNG, PDF ou PPTX).",
            )
            return

        per_page_value = self.page_layout.get().strip()
        try:
            per_page = int(float(per_page_value))
        except ValueError:
            messagebox.showerror(
                "Configuração inválida",
                "Informe 4 ou 6 gráficos por página para gerar o relatório.",
            )
            return

        if per_page not in {4, 6}:
            messagebox.showerror(
                "Configuração inválida",
                "A quantidade de gráficos por página deve ser 4 ou 6.",
            )
            return

        if should_export_pdf or should_export_ppt:
            try:
                logos = self._load_brand_assets()
            except FileNotFoundError as exc:
                messagebox.showerror("Erro", str(exc))
                return
        else:
            logos = {}

        kpis = self.load_kpi_formulas()
        if not kpis:
            messagebox.showerror("Erro", "Arquivo kpi_formulas.json não encontrado.")
            return

        filter_selections = self._collect_filter_selections()
        date_ini = self.date_ini.get().strip()
        date_fim = self.date_fim.get().strip()
        selected_hours = [self.hour_filter.get(i) for i in self.hour_filter.curselection()]
        tasklines_snapshot = list(self.tasklines)
        report_name_entry = self.kpi_title_entry.get().strip()
        followup_slug = self._normalize_report_name(report_name_entry)
        data_hoje = datetime.now().strftime("%Y%m%d")
        filter_slug = self._build_filter_slug(filter_selections)
        base_folder = self._compose_base_folder(followup_slug, data_hoje)
        file_basename = self._compose_file_basename(followup_slug, filter_slug, data_hoje)
        dpi_value = 400 if self.high_definition.get() else 100

        self._is_generating = True
        self._show_progress("Preparando geração...", 0)

        thread = threading.Thread(
            target=self._generate_report_thread,
            args=(
                self.df,
                filter_selections,
                (date_ini, date_fim),
                selected_hours,
                tasklines_snapshot,
                kpis,
                logos,
                followup_slug,
                data_hoje,
                file_basename,
                base_folder,
                per_page,
                bool(should_export_png),
                bool(should_export_pdf),
                bool(should_export_ppt),
                dpi_value,
                report_name_entry,
            ),
            daemon=True,
        )
        thread.start()

    def _generate_report_thread(
        self,
        df_source: pd.DataFrame | None,
        filter_selections: dict[str, list[str]],
        date_range: tuple[str | None, str | None],
        selected_hours: list[str],
        tasklines_snapshot: list[tuple[pd.Timestamp, str]],
        kpis: dict,
        logos: dict[str, str],
        followup_slug: str,
        data_hoje: str,
        base_file_name: str,
        base_folder: str,
        per_page: int,
        should_export_png: bool,
        should_export_pdf: bool,
        should_export_ppt: bool,
        dpi_value: int,
        report_name_entry: str,
    ) -> None:
        start_date, end_date = date_range
        try:
            df_filt = self.aplicar_filtros(
                df_source,
                selections=filter_selections,
                date_range=(start_date, end_date),
                selected_hours=selected_hours,
            )
            if "DATETIME" not in df_filt.columns:
                self._schedule_on_main(
                    self._handle_generation_failure,
                    "Erro",
                    "Coluna DATETIME não encontrada nos dados filtrados. Verifique os arquivos de entrada.",
                )
                return

            df_filt = df_filt.dropna(subset=["DATETIME"])
            if df_filt.empty:
                self._schedule_on_main(
                    self._handle_generation_failure,
                    "Sem dados",
                    "Nenhum dado encontrado com os filtros selecionados.",
                )
                return

            metadata = self._gather_report_metadata(
                df_filt,
                filter_selections=filter_selections,
                date_ini=start_date,
                date_fim=end_date,
                selected_hours=selected_hours,
                tasklines=tasklines_snapshot,
                report_name_entry=report_name_entry,
            )

            out_dir = os.path.join(os.getcwd(), base_folder)
            os.makedirs(out_dir, exist_ok=True)

            messages = self._run_report_generation(
                df_filt,
                metadata,
                kpis,
                should_export_png,
                should_export_pdf,
                should_export_ppt,
                out_dir,
                followup_slug,
                data_hoje,
                base_file_name,
                per_page,
                logos,
                dpi_value,
                tasklines=tasklines_snapshot,
            )
        except PermissionError:
            self._schedule_on_main(
                self._handle_generation_failure,
                "Erro ao salvar",
                "O arquivo está aberto em outro programa. Feche-o e tente novamente.",
            )
            return
        except Exception as exc:
            self._schedule_on_main(self._handle_generation_failure, "Erro", str(exc))
            return

        self._schedule_on_main(self._handle_generation_success, messages)

    def _schedule_on_main(self, callback, *args, **kwargs) -> None:
        def _runner():
            callback(*args, **kwargs)

        self.master.after(0, _runner)

    def _queue_progress(self, message: str, value: float) -> None:
        self._schedule_on_main(self._show_progress, message, value)

    def _queue_hide_progress(self) -> None:
        self._schedule_on_main(self._hide_progress)

    def _handle_generation_success(self, messages: list[str]) -> None:
        self._is_generating = False
        self._hide_progress()
        if messages:
            info_lines = "\n".join(messages)
            messagebox.showinfo("Concluído", f"Relatório gerado com sucesso!\n{info_lines}")
        else:
            messagebox.showinfo("Concluído", "Relatório gerado com sucesso!")

    def _handle_generation_failure(self, title: str, message: str) -> None:
        self._is_generating = False
        self._hide_progress()
        messagebox.showerror(title, message)

    def _run_report_generation(
        self,
        df_filt: pd.DataFrame,
        metadata: dict,
        kpis: dict,
        should_export_png: bool,
        should_export_pdf: bool,
        should_export_ppt: bool,
        out_dir: str,
        followup_slug: str,
        data_hoje: str,
        base_file_name: str,
        per_page: int,
        logos: dict[str, str],
        dpi_value: int,
        *,
        tasklines: list[tuple[pd.Timestamp, str]] | None = None,
    ) -> list[str]:
        charts: list[dict[str, object]] = []
        ppt_temp_paths: list[str] = []
        if tasklines is None:
            tasklines = list(self.tasklines)
        generated_messages: list[str] = []
        failures: list[str] = []
        log_written = False

        def ensure_failure_log() -> str | None:
            nonlocal log_written
            if not failures or log_written:
                return None
            error_log_path = os.path.join(out_dir, "kpi_erros.log")
            with open(error_log_path, "w", encoding="utf-8") as log_file:
                log_file.write("KPIs ignorados durante a geração do relatório:\n")
                for item in failures:
                    log_file.write(f"- {item}\n")
            log_written = True
            return error_log_path

        followup_suffix = f"_{followup_slug}" if followup_slug else ""
        png_dir = os.path.join(out_dir, "png") if should_export_png else out_dir
        if should_export_png:
            os.makedirs(png_dir, exist_ok=True)
        sanitized_columns: set[str] = set()
        chart_fig = plt.figure(figsize=(9.4, 5.6))
        chart_ax = chart_fig.add_subplot(111)
        try:
            total_kpis = max(len(kpis), 1)
            for processed_kpis, (kpi_name, info) in enumerate(kpis.items(), start=1):
                progress_pct = (processed_kpis / total_kpis) * 70
                self._queue_progress(
                    f"Gerando gráfico {processed_kpis}/{total_kpis}...",
                    progress_pct,
                )
                formula_text = info.get("formula", "")
                required_columns = self._extract_formula_columns(formula_text)
                missing_columns = sorted(col for col in required_columns if col not in df_filt.columns)
                if missing_columns:
                    failure_msg = f"{kpi_name}: colunas ausentes - {', '.join(missing_columns)}"
                    failures.append(failure_msg)
                    print(f"[KPI ignorado] {failure_msg}")
                    continue

                for column in required_columns:
                    if column in sanitized_columns:
                        continue
                    series = df_filt[column]
                    if is_numeric_dtype(series) and series.dtype == np.float32:
                        sanitized_columns.add(column)
                        continue
                    df_filt[column] = self._sanitize_numeric_series(series)
                    sanitized_columns.add(column)

                subset_columns = ["DATETIME"]
                if required_columns:
                    subset_columns.extend(sorted(required_columns))
                df_subset = df_filt.loc[:, subset_columns]
                grouped = df_subset.groupby("DATETIME", sort=True)

                replaced_formula = self._replace_source_refs(formula_text)
                try:
                    compiled_formula = compile(
                        replaced_formula,
                        f"<kpi:{self._slugify(kpi_name) or 'formula'}>",
                        "eval",
                    )
                except SyntaxError as exc:
                    failure_msg = f"{kpi_name}: fórmula inválida - {exc}"
                    failures.append(failure_msg)
                    print(f"[KPI ignorado] {failure_msg}")
                    continue

                try:
                    result_df = self._evaluate_formula(grouped, compiled_formula, kpi_name)
                except Exception as exc:
                    failure_msg = f"{kpi_name}: erro ao calcular - {exc}"
                    failures.append(failure_msg)
                    print(f"[KPI ignorado] {failure_msg}")
                    continue

                has_data = not result_df.empty
                if not has_data:
                    failure_msg = f"{kpi_name}: sem dados após o cálculo"
                    failures.append(failure_msg)
                    print(f"[KPI ignorado] {failure_msg}")

                is_percent = self._is_percentage_metric(formula_text)
                chart_type = info.get("chart_type", "line")

                chart_ax.clear()
                if has_data:
                    self._render_chart(
                        chart_ax,
                        result_df,
                        kpi_name,
                        is_percent,
                        chart_type,
                        style_scale=1.0,
                        tasklines=tasklines,
                    )
                else:
                    self._render_placeholder_chart(
                        chart_ax,
                        kpi_name,
                        style_scale=1.0,
                    )
                chart_fig.subplots_adjust(left=0.06, right=0.98, bottom=0.14, top=0.86)

                kpi_slug = self._slugify(kpi_name)
                nome_base = (
                    f"{base_file_name}_{kpi_slug}"
                    if base_file_name
                    else f"FollowUp_5G{followup_suffix}_{kpi_slug}_{data_hoje}"
                )
                img_path = os.path.join(png_dir, nome_base + ".png")
                chart_fig.savefig(img_path, dpi=dpi_value, bbox_inches="tight", pad_inches=0.25)

                ppt_image_path: str | None = None
                if should_export_ppt:
                    chart_ax.clear()
                    if has_data:
                        self._render_chart(
                            chart_ax,
                            result_df,
                            kpi_name,
                            is_percent,
                            chart_type,
                            style_scale=1.2,
                            tasklines=tasklines,
                        )
                    else:
                        self._render_placeholder_chart(
                            chart_ax,
                            kpi_name,
                            style_scale=1.2,
                        )
                    chart_fig.subplots_adjust(left=0.06, right=0.98, bottom=0.14, top=0.86)
                    ppt_image_path = os.path.join(out_dir, nome_base + "_ppt.png")
                    chart_fig.savefig(
                        ppt_image_path,
                        dpi=dpi_value,
                        bbox_inches="tight",
                        pad_inches=0.25,
                    )
                    ppt_temp_paths.append(ppt_image_path)

                charts.append(
                    {
                        "image_path": img_path,
                        "data": result_df if has_data else None,
                        "name": kpi_name,
                        "is_percent": is_percent,
                        "chart_type": chart_type,
                        "ppt_image_path": ppt_image_path,
                    }
                )

            chart_fig.clf()
            if not charts:
                log_path = ensure_failure_log()
                if log_path:
                    raise ValueError("Nenhum gráfico foi gerado. Consulte kpi_erros.log para detalhes.")
                raise ValueError("Nenhum gráfico foi gerado.")
            rows = 2
            cols = 2 if per_page == 4 else 3
            page_size = (11.69, 8.27)

            if should_export_pdf:
                if not logos:
                    raise FileNotFoundError(
                        "Arquivos de branding ausentes para montar o PDF."
                    )
                self._queue_progress("Montando PDF...", 75)
                pdf_filename = f"{base_file_name}.pdf" if base_file_name else f"FollowUp_5G{followup_suffix}_{data_hoje}.pdf"
                pdf_path = os.path.join(out_dir, pdf_filename)
                style_scale = 0.8 if per_page == 4 else 0.62
                with PdfPages(pdf_path) as pdf:
                    self._build_pdf_cover(pdf, metadata, logos, page_size, dpi_value)
                    for chunk in self._chunked(charts, per_page):
                        page_fig, axs = plt.subplots(rows, cols, figsize=page_size)
                        axs_array = np.array(axs).reshape(rows * cols)
                        for ax in axs_array:
                            ax.clear()
                            ax.axis("off")
                        for ax, chart in zip(axs_array, chunk):
                            ax.axis("on")
                            chart_data = chart.get("data")
                            if chart_data is None:
                                self._render_placeholder_chart(
                                    ax,
                                    chart["name"],
                                    style_scale=style_scale,
                                )
                            else:
                                self._render_chart(
                                    ax,
                                    chart_data,
                                    chart["name"],
                                    bool(chart["is_percent"]),
                                    chart.get("chart_type", "line"),
                                    style_scale=style_scale,
                                    tasklines=tasklines,
                                )
                        page_fig.subplots_adjust(
                            left=0.08,
                            right=0.97,
                            bottom=0.12,
                            top=0.78,
                            wspace=0.38,
                            hspace=0.58,
                        )
                        self._add_pdf_header(page_fig, metadata, logos)
                        pdf.savefig(page_fig, dpi=dpi_value, bbox_inches="tight")
                        plt.close(page_fig)
                generated_messages.append(f"PDF: {os.path.basename(pdf_path)}")

            if should_export_ppt:
                if not logos:
                    raise FileNotFoundError(
                        "Arquivos de branding ausentes para montar o PPTX."
                    )
                self._queue_progress("Montando PPTX...", 85)
                ppt_filename = f"{base_file_name}.pptx" if base_file_name else f"FollowUp_5G{followup_suffix}_{data_hoje}.pptx"
                ppt_path = os.path.join(out_dir, ppt_filename)
                prs = Presentation()
                prs.slide_width = Inches(13.33)
                prs.slide_height = Inches(7.5)
                blank = prs.slide_layouts[6]

                self._add_ppt_cover(prs, metadata, logos)

                margin_w = Inches(0.4)
                margin_h = Inches(0.4)
                spacing_w = Inches(0.2)
                spacing_h = Inches(0.2)

                usable_width = prs.slide_width - 2 * margin_w - (cols - 1) * spacing_w
                usable_height = prs.slide_height - 2 * margin_h - (rows - 1) * spacing_h
                img_width = usable_width / cols
                img_height = usable_height / rows

                for chunk in self._chunked(charts, per_page):
                    slide = prs.slides.add_slide(blank)
                    for idx, chart in enumerate(chunk):
                        row = idx // cols
                        col = idx % cols
                        left = margin_w + col * (img_width + spacing_w)
                        top = margin_h + row * (img_height + spacing_h)
                        img_path = chart.get("ppt_image_path") or chart["image_path"]
                        slide.shapes.add_picture(
                            img_path, left, top, width=img_width, height=img_height
                        )
                prs.save(ppt_path)
                generated_messages.append(f"PPTX: {os.path.basename(ppt_path)}")

            for chart in charts:
                chart.pop("data", None)

            if should_export_png:
                generated_messages.append("PNGs: arquivos salvos na pasta 'png'.")

            log_path = ensure_failure_log()
            if log_path:
                generated_messages.append(
                    f"{len(failures)} KPI(s) foram ignorados. Consulte kpi_erros.log para detalhes."
                )

            self._queue_progress("Finalizando...", 100)
            return generated_messages

        finally:
            plt.close(chart_fig)
            if not should_export_png:
                for chart in charts:
                    img_path = chart.get("image_path")
                    if img_path and os.path.exists(img_path):
                        try:
                            os.remove(img_path)
                        except OSError:
                            pass
            for ppt_path in ppt_temp_paths:
                if ppt_path and os.path.exists(ppt_path):
                    try:
                        os.remove(ppt_path)
                    except OSError:
                        pass
            charts.clear()
            plt.close("all")
            gc.collect()
        raise ValueError("Nenhum gráfico foi gerado.")

    def _show_progress(self, message: str, value: float):
        self.progress_frame.pack(fill="x", pady=(5, 0))
        self.progress_label.config(text=message)
        self.progress_var.set(max(0.0, min(100.0, value)))
        self.master.update_idletasks()

    def _hide_progress(self):
        self.progress_var.set(0.0)
        self.progress_label.config(text="")
        self.progress_frame.pack_forget()


if __name__ == "__main__":
    root = tk.Tk()
    app = FollowUpApp(root)
    root.mainloop()
